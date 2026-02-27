#!/usr/bin/env python3
"""
Quantum Shield - Pitch Deck Generator
Generates a .pptx file importable to Google Slides.

Usage:
  python3 docs/pitch/generate_slides.py

Output:
  docs/pitch/quantum-shield-pitch.pptx

Requirements:
  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
HINOMARU_RED = RGBColor(0xC4, 0x1E, 0x3A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF8, 0xFC)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "quantum-shield-pitch.pptx")
SCREENSHOT_DIR = os.path.join(os.path.dirname(__file__), "screenshots")

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_NAVY):
    """Fill slide background with color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=0):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, width, height, lines, size=14, color=WHITE, line_spacing=1.5):
    """Add multi-line text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(size * (line_spacing - 1))
    return txBox


# ===== SLIDE 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide, DARK_NAVY)

# Red accent bar at top
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

# Title
add_text(slide, 1.5, 1.5, 10, 1.2, "QUANTUM SHIELD", 54, HINOMARU_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 2.7, 10, 0.8, "The First Quantum-Safe Asset Protection Protocol", 24, GOLD, align=PP_ALIGN.CENTER)

# Divider line
add_shape(slide, 4.5, 3.7, 4.3, 0.02, GOLD)

# Subtitle
add_text(slide, 1.5, 4.0, 10, 0.6, "Protecting $2.5T in smart contract assets", 20, WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.5, 10, 0.6, "from the quantum threat — before it's too late.", 20, WHITE, align=PP_ALIGN.CENTER)

# Bottom info
add_text(slide, 1.5, 6.0, 5, 0.5, "Seed Round  |  February 2026  |  Confidential", 14, MED_GRAY, align=PP_ALIGN.CENTER)

# Red bottom bar
add_shape(slide, 0, 7.42, 13.333, 0.08, HINOMARU_RED)


# ===== SLIDE 2: Problem =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "The Quantum Clock is Ticking", 36, DARK_NAVY, bold=True)

# Stat boxes
stats = [
    ("2028", "US Federal\nPQ Deadline"),
    ("$2.5T", "Smart Contract\nAssets at Risk"),
    ("NOW", "\"Harvest Now,\nDecrypt Later\""),
    ("0", "Production PQ\nProtection on ETH"),
]

for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 3.1
    box = add_shape(slide, x, 1.5, 2.8, 1.8, LIGHT_GRAY)
    add_text(slide, x + 0.2, 1.6, 2.4, 0.9, num, 40, HINOMARU_RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 2.4, 2.4, 0.7, label, 13, MED_GRAY, align=PP_ALIGN.CENTER)

# Key points
points = [
    "NIST finalized post-quantum standards (FIPS 204, 205) in August 2024",
    "US Executive Order 14110: Federal agencies must transition by 2028",
    "EU DORA/MiCA 2.0: Quantum-safe requirements for financial infrastructure",
    "Vitalik Buterin: \"Realistic worst case for quantum threat is 2028\"",
    "PsiQuantum raised $1B (BlackRock, Nvidia) in Sep 2025 — quantum is real",
]
add_multiline(slide, 0.8, 3.8, 11.5, 3, [f"•  {p}" for p in points], 16, DARK_TEXT, 1.8)


# ===== SLIDE 3: Market =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Market: Multi-Billion Dollar PQC Wave", 36, DARK_NAVY, bold=True)

# TAM/SAM/SOM
markets = [
    ("TAM", "$2.8 - $4.6B", "PQC Market by 2030", "MarketsandMarkets, Mordor Intelligence\n37-46% CAGR", HINOMARU_RED),
    ("SAM", "$400 - $900M", "Crypto/DeFi PQ Infra", "15-20% of TAM\nEU DORA + DeFi protocols", GOLD),
    ("SOM", "$50 - $150M", "ETH-native PQ (2026-28)", "Early adopters\nInstitutional custody", DARK_NAVY),
]

for i, (label, amount, desc, detail, color) in enumerate(markets):
    y = 1.4 + i * 1.7
    # Left label
    lbl = add_shape(slide, 0.8, y, 1.2, 1.3, color)
    add_text(slide, 0.8, y + 0.3, 1.2, 0.6, label, 24, WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Content
    add_text(slide, 2.3, y, 4, 0.6, amount, 28, DARK_NAVY, bold=True)
    add_text(slide, 2.3, y + 0.5, 4, 0.4, desc, 16, MED_GRAY)
    add_text(slide, 2.3, y + 0.9, 4, 0.5, detail, 11, MED_GRAY)

# Right side: Window
add_shape(slide, 7.5, 1.4, 5, 5.2, DARK_NAVY)
add_text(slide, 7.8, 1.7, 4.4, 0.6, "CRITICAL WINDOW", 22, GOLD, bold=True, align=PP_ALIGN.CENTER)
window_lines = [
    "Ethereum PQ Upgrade: ~2028",
    "Our head start: 2 years",
    "",
    "QANplatform raised $17.1M",
    "→ PQ × Crypto has VC appetite",
    "",
    "PsiQuantum raised $1B",
    "→ Quantum attracting serious capital",
    "",
    "US Federal budget: $7.1B",
    "→ for PQ migration",
]
add_multiline(slide, 7.8, 2.5, 4.4, 3.5, window_lines, 14, WHITE, 1.5)


# ===== SLIDE 4: Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "3-Layer Quantum-Safe Architecture", 36, DARK_NAVY, bold=True)

# Three layers
layers = [
    ("L1 Vault", "SPHINCS+ (FIPS 205)", "On-chain, immutable vault\nEthereum Mainnet", HINOMARU_RED),
    ("L3 Aegis", "Dilithium (FIPS 204)", "Off-chain BFT consensus\n93% gas savings", GOLD),
    ("Prover Pool", "VRF + Quadratic Slashing", "Decentralized operators\nEconomic security bond", DARK_NAVY),
]

for i, (name, algo, desc, color) in enumerate(layers):
    y = 1.4 + i * 1.6
    add_shape(slide, 0.8, y, 0.15, 1.3, color)
    add_text(slide, 1.2, y, 3, 0.5, name, 24, DARK_NAVY, bold=True)
    add_text(slide, 1.2, y + 0.45, 3, 0.4, algo, 14, color, bold=True)
    add_text(slide, 1.2, y + 0.8, 4, 0.6, desc, 13, MED_GRAY)

# Right side: User flow + Differentiators
add_shape(slide, 7, 1.4, 5.5, 2.5, LIGHT_GRAY)
add_text(slide, 7.3, 1.5, 5, 0.5, "User Flow", 18, DARK_NAVY, bold=True)
add_text(slide, 7.3, 2.0, 5, 0.5, "Lock (2s) → 24h TimeLock → Auto-Claim → Unlock", 15, DARK_TEXT)
add_text(slide, 7.3, 2.7, 5, 0.4, "Zero manual action needed for unlock", 13, MED_GRAY)
add_text(slide, 7.3, 3.1, 5, 0.4, "Emergency Path: 7d failsafe, guaranteed recovery", 13, MED_GRAY)

# Differentiator pills
pills = ["Dual NIST", "Quadratic Slashing", "Auto-Claim", "24h TimeLock", "Emergency Path"]
for i, pill in enumerate(pills):
    x = 7 + (i % 3) * 1.8
    y = 4.3 + (i // 3) * 0.5
    p = add_shape(slide, x, y, 1.6, 0.35, HINOMARU_RED)
    add_text(slide, x, y + 0.02, 1.6, 0.3, pill, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)


# ===== SLIDE 5: Competitive Position =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Competitive Position: Only Dual-NIST Protocol", 36, DARK_NAVY, bold=True)

# Comparison table
headers = ["Feature", "QRL", "QANplatform", "ETH PQ", "Solana", "QS ★"]
rows = [
    ["Dilithium (FIPS 204)", "—", "✓", "~2028", "✓", "✓"],
    ["SPHINCS+ (FIPS 205)", "✓", "—", "~2028", "—", "✓"],
    ["Economic Security", "—", "—", "—", "—", "✓"],
    ["Auto-Claim UX", "—", "—", "—", "—", "✓"],
    ["Full App Ecosystem", "1 app", "1 app", "—", "—", "9 apps"],
    ["Ethereum-native", "Own L1", "Own L1", "Core", "Own L1", "✓"],
    ["Status", "Mainnet", "Pilot", "~2028", "Testnet", "Testnet"],
    ["Funding", "$0 (ICO '17)", "$17.1M", "∞", "∞", "Seeking"],
]

# Table header
for j, h in enumerate(headers):
    x = 0.8 + j * 2.05
    c = HINOMARU_RED if j == len(headers) - 1 else DARK_NAVY
    add_shape(slide, x, 1.4, 2.0, 0.45, c)
    add_text(slide, x, 1.42, 2.0, 0.4, h, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)

# Table rows
for i, row in enumerate(rows):
    y = 1.9 + i * 0.5
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = 0.8 + j * 2.05
        if j == len(row) - 1:
            # Highlight QS column
            add_shape(slide, x, y, 2.0, 0.45, RGBColor(0xFF, 0xF5, 0xF5))
        elif j == 0:
            pass  # No background for label
        add_text(slide, x, y + 0.02, 2.0, 0.4, cell, 11 if j == 0 else 12,
                 DARK_TEXT if j == 0 else (HINOMARU_RED if j == len(row) - 1 else DARK_TEXT),
                 bold=(j == len(row) - 1 and cell == "✓"),
                 align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


# ===== SLIDE 6: Quadratic Slashing =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Novel: Quadratic Slashing", 36, GOLD, bold=True)
add_text(slide, 0.8, 1.1, 12, 0.5, "penalty = N² × 10% of stake", 22, WHITE)
add_text(slide, 0.8, 1.6, 12, 0.4, "Making collusion exponentially more expensive than honesty", 16, MED_GRAY)

# Penalty visualization
penalties = [
    ("1 cheater", "10%", "90% remains", RGBColor(0x28, 0xA7, 0x45)),
    ("2 colluders", "40%", "60% remains", GOLD),
    ("3 colluders", "90%", "10% remains", HINOMARU_RED),
    ("4+ colluders", "100%+", "0% remains", RGBColor(0xFF, 0x00, 0x00)),
]

for i, (label, penalty, remaining, color) in enumerate(penalties):
    y = 2.4 + i * 1.1
    # Bar representing loss
    bar_width = float(penalty.replace('%', '').replace('+', '')) / 100 * 8
    if bar_width > 8:
        bar_width = 8
    add_shape(slide, 0.8, y, bar_width, 0.7, color)
    add_text(slide, 0.8, y + 0.1, bar_width if bar_width > 1 else 2, 0.5, f"{label}: {penalty} lost", 16, WHITE, bold=True)
    add_text(slide, 9.5, y + 0.15, 3, 0.4, remaining, 14, color, bold=True)

# Right side: Dollar example
add_shape(slide, 9, 2.4, 3.8, 3.5, RGBColor(0x22, 0x22, 0x40))
add_text(slide, 9.2, 2.5, 3.4, 0.4, "With $400K stake/Prover:", 14, GOLD, bold=True)
example_lines = [
    "1 cheater: $40K lost",
    "2 colluders: $320K total",
    "3 colluders: $1.08M total",
    "",
    "60% of slashed funds →",
    "Challenger reward",
    "= active fraud detection"
]
add_multiline(slide, 9.2, 3.0, 3.4, 2.5, example_lines, 13, WHITE, 1.4)


# ===== SLIDE 7: Product =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Product: 9 Apps, 175 Screens, Ready to Ship", 36, DARK_NAVY, bold=True)
add_text(slide, 0.8, 1.0, 12, 0.4, "Full ecosystem — not just a smart contract", 18, MED_GRAY)

# App tiers
tiers = [
    ("Consumer", HINOMARU_RED, [
        ("Consumer App", "Lock / Unlock / Emergency"),
        ("Explorer", "Public tx history, risk scores"),
        ("Token Hub", "veQS staking, rewards"),
    ]),
    ("Governance", GOLD, [
        ("Governance Portal", "Proposals, veQS voting"),
        ("QS Hub", "Treasury, allocation tracking"),
    ]),
    ("Institutional", DARK_NAVY, [
        ("Enterprise Admin", "Multi-sig custody"),
        ("QS Admin", "Foundation operations"),
    ]),
    ("Operators", MED_GRAY, [
        ("Prover Portal", "Node ops, signing queue"),
        ("Observer Portal", "Monitoring, challenges"),
    ]),
]

for i, (tier_name, color, apps) in enumerate(tiers):
    x = 0.5 + i * 3.2
    add_shape(slide, x, 1.6, 3.0, 0.4, color)
    add_text(slide, x, 1.62, 3.0, 0.35, f"  {tier_name}", 14, WHITE, bold=True)

    for j, (app_name, app_desc) in enumerate(apps):
        y = 2.2 + j * 1.0
        add_shape(slide, x, y, 3.0, 0.85, LIGHT_GRAY)
        add_text(slide, x + 0.15, y + 0.05, 2.7, 0.35, app_name, 13, DARK_NAVY, bold=True)
        add_text(slide, x + 0.15, y + 0.4, 2.7, 0.4, app_desc, 10, MED_GRAY)

# Stats bar at bottom
stats_bottom = [
    ("375", "Components"),
    ("202", "API Endpoints"),
    ("144", "E2E Tests"),
    ("251", "Pages"),
    ("Live", "Sepolia Testnet"),
]
add_shape(slide, 0.5, 5.8, 12.3, 1.0, DARK_NAVY)
for i, (val, label) in enumerate(stats_bottom):
    x = 1.0 + i * 2.4
    add_text(slide, x, 5.85, 2.0, 0.5, val, 28, GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 6.3, 2.0, 0.3, label, 11, WHITE, align=PP_ALIGN.CENTER)

# Screenshot placeholder
add_text(slide, 0.8, 7.0, 12, 0.3,
         "→ Replace with actual screenshots from capture-pitch-screenshots.ts",
         10, MED_GRAY)


# ===== SLIDE 8: Traction =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Traction: 92% Complete", 36, DARK_NAVY, bold=True)

# Progress bar
add_shape(slide, 0.8, 1.3, 11.7, 0.5, LIGHT_GRAY)
add_shape(slide, 0.8, 1.3, 11.7 * 0.92, 0.5, HINOMARU_RED)
add_text(slide, 0.8, 1.32, 11.7, 0.45, "  92%", 20, WHITE, bold=True)

# Done items
done_items = [
    "Core Protocol (SEQUENCES v3.0) — 9 complete sequences",
    "L1 Vault on Sepolia (deployed, verifiable)",
    "9 Frontend Applications (React/Next.js)",
    "202 Backend API Endpoints (Rust/Axum)",
    "PostgreSQL + Redis Storage Layer",
    "144 E2E Test Files (Playwright)",
]
add_text(slide, 0.8, 2.1, 5.5, 0.4, "Done", 20, RGBColor(0x28, 0xA7, 0x45), bold=True)
add_multiline(slide, 0.8, 2.5, 5.5, 3, [f"✅  {item}" for item in done_items], 13, DARK_TEXT, 1.6)

# Remaining
remaining = [
    "L3 Environment Configuration",
    "Mock Data Cleanup (38 locations)",
    "Final E2E Execution",
]
add_text(slide, 7, 2.1, 5.5, 0.4, "Remaining (~15 days)", 20, GOLD, bold=True)
add_multiline(slide, 7, 2.5, 5.5, 2, [f"🔜  {item}" for item in remaining], 13, DARK_TEXT, 1.6)

# Milestones
add_shape(slide, 0.8, 5.2, 11.7, 2.0, DARK_NAVY)
milestones = [
    ("2025 Q3", "Architecture Design"),
    ("2025 Q4", "Core Protocol + Smart Contracts"),
    ("2026 Q1 ← NOW", "9 Apps + Backend + Testnet"),
    ("2026 Q2", "Audit → Mainnet"),
]
for i, (quarter, desc) in enumerate(milestones):
    x = 1 + i * 3
    is_now = "NOW" in quarter
    c = GOLD if is_now else MED_GRAY
    add_shape(slide, x, 5.5, 2.5, 0.05, c)
    add_text(slide, x, 5.6, 2.5, 0.4, quarter, 14, GOLD if is_now else WHITE, bold=is_now, align=PP_ALIGN.CENTER)
    add_text(slide, x, 6.0, 2.5, 0.4, desc, 11, WHITE if is_now else MED_GRAY, align=PP_ALIGN.CENTER)


# ===== SLIDE 9: Business Model =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Business Model & Unit Economics", 36, DARK_NAVY, bold=True)

# Transaction economics
add_shape(slide, 0.8, 1.4, 5.5, 3.0, LIGHT_GRAY)
add_text(slide, 1.0, 1.5, 5, 0.4, "Transaction Economics", 20, DARK_NAVY, bold=True)
tx_lines = [
    "Lock:    ~$7   (135K gas)",
    "Unlock:  ~$27  (490K gas, Treasury-paid)",
    "Total:   ~$34  per round-trip",
    "",
    "vs. On-chain PQ: $40-100+",
    "→ 30-70% cheaper via L1/L3 split",
]
add_multiline(slide, 1.0, 2.0, 5, 2, tx_lines, 14, DARK_TEXT, 1.4)

# Revenue streams
add_shape(slide, 7, 1.4, 5.5, 3.0, LIGHT_GRAY)
add_text(slide, 7.2, 1.5, 5, 0.4, "Revenue Streams", 20, DARK_NAVY, bold=True)
rev_lines = [
    "•  Protocol Fees (0.1-0.5% of lock amount)",
    "•  Enterprise Licensing (multi-org features)",
    "•  Prover Rewards (consensus participation)",
    "•  Treasury Yield (protocol-owned liquidity)",
]
add_multiline(slide, 7.2, 2.0, 5, 2, rev_lines, 14, DARK_TEXT, 1.5)

# Token distribution
add_text(slide, 0.8, 4.7, 12, 0.4, "veQS Token Distribution", 20, DARK_NAVY, bold=True)

dist = [
    ("50%", "veQS Holders\n(Governance)", HINOMARU_RED),
    ("30%", "Provers\n(Security)", GOLD),
    ("10%", "Observers\n(Monitoring)", DARK_NAVY),
    ("10%", "Treasury\n(Development)", MED_GRAY),
]

# Distribution bar
bar_start = 0.8
for pct_str, label, color in dist:
    pct = int(pct_str.replace('%', ''))
    width = 11.7 * pct / 100
    add_shape(slide, bar_start, 5.2, width, 0.6, color)
    add_text(slide, bar_start, 5.25, width, 0.5, f" {pct_str}", 16, WHITE, bold=True)
    add_text(slide, bar_start, 5.9, width, 0.6, label, 10, MED_GRAY, align=PP_ALIGN.CENTER)
    bar_start += width


# ===== SLIDE 10: Roadmap =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.4, 12, 0.7, "Roadmap", 36, GOLD, bold=True)

quarters_2026 = [
    ("Q1 2026", "NOW", [
        "Sepolia Testnet Public",
        "9 Apps Functional",
        "L3 Config + Mock Cleanup",
    ]),
    ("Q2 2026", "", [
        "Security Audit",
        "Pilot: 4-8 Trusted Provers",
        "Enterprise Beta Program",
    ]),
    ("Q3 2026", "", [
        "Ethereum Mainnet Launch",
        "veQS Governance Live",
        "Multi-asset (ERC-20)",
    ]),
    ("Q4 2026", "", [
        "Decentralized Prover Onboarding",
        "DeFi Partnerships",
        "Cross-chain Expansion",
    ]),
]

for i, (q, badge, items) in enumerate(quarters_2026):
    x = 0.5 + i * 3.2
    is_now = badge == "NOW"
    box_color = HINOMARU_RED if is_now else RGBColor(0x22, 0x22, 0x40)
    add_shape(slide, x, 1.4, 3.0, 3.5, box_color)
    add_text(slide, x, 1.5, 3.0, 0.5, q, 20, GOLD if is_now else WHITE, bold=True, align=PP_ALIGN.CENTER)
    if badge:
        add_shape(slide, x + 1.0, 1.45, 1.0, 0.35, GOLD)
        add_text(slide, x + 1.0, 1.47, 1.0, 0.3, badge, 10, DARK_NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_multiline(slide, x + 0.2, 2.2, 2.6, 2.5, [f"•  {item}" for item in items], 12, WHITE, 1.5)

# 2027
add_shape(slide, 0.5, 5.3, 12.3, 1.5, RGBColor(0x22, 0x22, 0x40))
add_text(slide, 0.8, 5.4, 3, 0.5, "2027", 22, GOLD, bold=True)
goals_2027 = [
    "Hardware Wallet Integration",
    "Enterprise SaaS Launch",
    "Target: $500M+ TVL, 10M+ Locks",
]
add_multiline(slide, 3.5, 5.4, 8, 1.2, [f"•  {g}" for g in goals_2027], 14, WHITE, 1.4)


# ===== SLIDE 11: The Ask =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

add_text(slide, 0.8, 0.3, 12, 0.5, "Seed Round", 18, MED_GRAY)
add_text(slide, 0.8, 0.7, 12, 1.0, "$3 — $5M", 60, GOLD, bold=True)

# Use of Funds
funds = [
    ("40%", "Engineering & Audit", "Security audit ($200-500K)\nL1/L3 production optimization\nChief Cryptographer hire", HINOMARU_RED),
    ("35%", "Go-to-Market", "Enterprise Sales (2 hires)\nDeFi Partnership Development\nCommunity Building", GOLD),
    ("15%", "Operations", "L3 Infrastructure (12 months)\nProver Pilot Program", RGBColor(0x55, 0x55, 0x88)),
    ("10%", "Legal & Compliance", "Token Legal Structure\nRegulatory Compliance", MED_GRAY),
]

for i, (pct, title, details, color) in enumerate(funds):
    x = 0.5 + i * 3.2
    add_shape(slide, x, 2.0, 3.0, 3.0, RGBColor(0x22, 0x22, 0x40))
    add_shape(slide, x, 2.0, 3.0, 0.5, color)
    add_text(slide, x, 2.02, 1.5, 0.45, f" {pct}", 24, WHITE, bold=True)
    add_text(slide, x + 1.2, 2.05, 1.6, 0.4, title, 12, WHITE, bold=True)
    add_multiline(slide, x + 0.2, 2.7, 2.6, 2, details.split('\n'), 11, MED_GRAY, 1.5)

# 18-month targets
add_shape(slide, 0.5, 5.5, 12.3, 1.5, RGBColor(0x22, 0x22, 0x40))
add_text(slide, 0.8, 5.6, 12, 0.4, "18-Month Targets", 18, GOLD, bold=True)
targets = [
    ("$500M+", "TVL"),
    ("4+", "Institutional\nClients"),
    ("8+", "Active\nProvers"),
    ("Mainnet", "Q3 2026\nLaunch"),
]
for i, (val, label) in enumerate(targets):
    x = 1.5 + i * 2.8
    add_text(slide, x, 5.9, 2.2, 0.5, val, 28, GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 6.4, 2.2, 0.5, label, 11, WHITE, align=PP_ALIGN.CENTER)


# ===== SLIDE 12: Why Now, Why Us =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape(slide, 0, 0, 13.333, 0.08, HINOMARU_RED)

# Left: Why Now
add_text(slide, 0.8, 0.5, 5.5, 0.5, "Why Now", 30, GOLD, bold=True)
why_now = [
    "NIST standardized PQ algorithms (Aug 2024)",
    "US/EU mandates driving adoption (2028)",
    "Ethereum PQ upgrade: 2028 → 2yr window",
    "\"Harvest Now, Decrypt Later\" is happening",
]
add_multiline(slide, 0.8, 1.2, 5.5, 2.5, [f"•  {w}" for w in why_now], 15, WHITE, 1.7)

# Right: Why Us
add_text(slide, 7, 0.5, 5.5, 0.5, "Why Us", 30, GOLD, bold=True)
why_us = [
    "Only dual-NIST protocol in crypto",
    "Full product (not just whitepaper)",
    "92% built, Sepolia live, demo-ready",
    "Novel economic security (Quad Slashing)",
]
add_multiline(slide, 7, 1.2, 5.5, 2.5, [f"•  {w}" for w in why_us], 15, WHITE, 1.7)

# Center quote
add_shape(slide, 2, 4.0, 9.3, 2.0, RGBColor(0x22, 0x22, 0x40))
add_shape(slide, 2, 4.0, 0.1, 2.0, GOLD)
quote_lines = [
    "In 2 years, every major protocol will need",
    "quantum-safe infrastructure.",
    "",
    "We're building it today.",
]
add_multiline(slide, 2.5, 4.2, 8.5, 1.5, quote_lines, 22, WHITE, 1.3)

# Contact
add_text(slide, 0.8, 6.5, 12, 0.5, "[Your Name]  |  [Email]  |  [Twitter/X]", 16, MED_GRAY, align=PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, 0, 7.42, 13.333, 0.08, HINOMARU_RED)


# ===== SAVE =====
prs.save(OUTPUT_PATH)
print(f"✅ Saved: {OUTPUT_PATH}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Format: 16:9 Widescreen")
print(f"   Import to Google Slides: File → Import slides → Upload this .pptx")
